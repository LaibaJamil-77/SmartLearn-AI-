import os
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from pptx import Presentation
import pdfplumber
from langchain.text_splitter import RecursiveCharacterTextSplitter
import faiss
import numpy as np
from datetime import datetime
import uuid
from langchain_openai import OpenAIEmbeddings

# Load environment variables
load_dotenv()

# Initialize OpenAI Embeddings
embeddings = OpenAIEmbeddings(
    model="text-embedding-ada-002",
)

# Initialize Groq LLM
llm = ChatGroq(
    model="llama3-70b-8192",
    temperature=0,
)

def extract_text_from_file(file_path):
    """Extract text from PDF or TXT files."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File '{file_path}' not found.")
    
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        else:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: PDF, TXT.")
    except Exception as e:
        raise ValueError(f"Failed to read file '{file_path}': {e}")

def split_text_into_chunks(text):
    """Split text into smaller chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    return text_splitter.split_text(text)

def generate_embeddings(chunks):
    """Generate embeddings for text chunks using OpenAI Embeddings."""
    return np.array(embeddings.embed_documents(chunks))

def create_faiss_index(embeddings):
    """Create FAISS index for vector search."""
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    return index

def retrieve_relevant_chunks(query, index, chunks, k=3):
    """Retrieve relevant text chunks for a query using OpenAI Embeddings."""
    query_embedding = np.array([embeddings.embed_query(query)])
    distances, indices = index.search(query_embedding, k)
    return [chunks[i] for i in indices[0]]

def generate_extra_explanation(content, context="topic", lang_choice='english'):
    """Generate teacher-like explanation for slide or document content."""
    if context == "topic":
        prompt = ChatPromptTemplate.from_messages(
            [
                (
                   "system",
                    f"""You are a warm, enthusiastic, and experienced teacher explaining a topic to students in a classroom. Start with a friendly greeting, like you're addressing the class in person (e.g., 'As-salamu Alaikum, meri pyari class!' or 'Hello, my awesome students!'). Then, explain the **content** in a natural, conversational way, as if you're teaching face-to-face. 
                    
                    - Use **simple, relatable language** that feels human and engaging.
                    - Include **real-life examples, analogies, or stories** to make the concept fun and easy to understand.
                    - Structure your explanation like this:
                        - **Greeting**: Start with a warm, friendly greeting to set the tone.
                        - **Simple Explanation**: Break down the concept in a way that feels like a teacher talking to students, avoiding technical jargon.
                        - **Real-World Example or Analogy**: Share a relatable, everyday example or a fun analogy to connect the concept to real life.
                        - **Step-by-Step Breakdown**: Explain the concept in clear, numbered steps, like a teacher guiding students through a process.
                        - **Common Mistakes**: Point out 1-2 common mistakes if students might make and how to avoid them, in a kind and encouraging way.
                    - Avoid using words like 'slide,' 'slide 1,',test,title, or 'slide 2'—treat the content as a topic you're teaching naturally.
                    - Keep the tone warm, encouraging, and conversational, like a favorite teacher who makes learning fun.
                    - Respond in {'Roman Urdu' if lang_choice == 'roman_urdu' else 'English'}, using natural, everyday phrases that students would hear in a classroom.
                    **Content to Explain:** {{content}}"""
                ),
            ]
        )
    elif context == "document":
        prompt = ChatPromptTemplate.from_messages(
            [
                (
                    "system",
                    f"""You are a warm, enthusiastic, and experienced teacher explaining a document to students in a classroom. Start with a friendly greeting, like you're addressing the class in person (e.g., 'As-salamu Alaikum, meri pyari class!' or 'Hello, my awesome students!'). Then, explain the **content** in a natural, conversational way, as if you're teaching face-to-face.
                    
                    - Use **simple, relatable language** that feels human and engaging.
                    - Include **real-life examples, analogies, or stories** to make the concept fun and easy to understand.
                    - Structure your explanation like this:
                        - **Greeting**: Start with a warm, friendly greeting to set the tone.
                        - **Simple Explanation**: Break down the concept in a way that feels like a teacher talking to students, avoiding technical jargon.
                        - **Real-World Example or Analogy**: Share a relatable, everyday example or a fun analogy to connect the concept to real life.
                        - **Step-by-Step Breakdown**: Explain the concept in clear, numbered steps, like a teacher guiding students through a process.
                        - **Common Mistakes**: Point out 1-2 common mistakes if students might make and how to avoid them, in a kind and encouraging way.
                    - Avoid using words like 'document' or 'content' in the explanation—treat it as a topic you're teaching naturally.
                    - Keep the tone warm, encouraging, and conversational, like a favorite teacher who makes learning fun.
                    - Respond in {'Roman Urdu' if lang_choice == 'roman_urdu' else 'English'}, using natural, everyday phrases that students would hear in a classroom.
                    **Content to Explain:** {{content}}"""
                ),
            ]
        )
    
    explanation_chain = prompt | llm
    try:
        return explanation_chain.invoke({"content": content}).content
    except Exception as e:
        print(f"Error generating explanation: {e}")
        return "Unable to generate explanation due to an error."

def generate_quiz(slides, context_type, lang_choice='english'):
    """Generate a quiz with 15 MCQs based on slide content."""
    slides_text = "\n".join([slide.strip() for slide in slides if slide.strip()])
    
    system_prompt = f"""You are an expert teacher tasked with creating a quiz with 15 multiple-choice questions (MCQs) based on the provided content. 
- Each question should have 4 answer options (A, B, C, D).
- Do NOT include the correct answer.
- Do NOT mention which option is correct.
- Just provide the question and the 4 options.
- Questions should cover key points from the content and vary in difficulty (easy, medium, hard).
- Format the quiz in Markdown, with each question starting with '# Question [number]', followed by the question and options.
- Respond in {'Roman Urdu' if lang_choice == 'roman_urdu' else 'English'}, using clear, student-friendly language.
- Ensure questions are relevant to the {'topic' if context_type == 'topic' else 'document'} content.
Content: {slides_text}"""

    
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
    ])
    chain = prompt | llm
    quiz_content = chain.invoke({"slides_text": slides_text}).content
    
    # Save quiz to a Markdown file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    quiz_filename = f"quiz_{context_type}_{timestamp}.md"
    with open(quiz_filename, 'w', encoding='utf-8') as f:
        f.write(f"# Quiz: {context_type.capitalize()}\n\n")
        f.write(quiz_content)
    
    print(f"Quiz saved as '{quiz_filename}'.")
    return quiz_content, quiz_filename

def generate_slides_from_topic(topic, instructions=None, num_slides=10, lang_choice='english'):
    """Generate presentation slides from a topic with clear, structured, teacher-style explanations."""
    system_prompt = f"""You are a helpful assistant that generates a {num_slides}-slide presentation on the topic: {{topic}}.
Each slide should be written in a clear, structured, and informative way, like a teacher explaining the content.
Start with an introduction and proceed logically through the topic, ending with a conclusion or summary.
Do not include headings like 'image' or mention any visuals, as images are not supported.
{instructions or 'Use standard presentation format with bullet points or concise paragraphs.'}
Each slide should begin with 'Slide [number]:' followed by the content in proper format."""
    # you can now pass this system_prompt to your LLM call

    prompt = ChatPromptTemplate.from_messages([("system", system_prompt)])
    chain = prompt | llm
    result = chain.invoke({"topic": topic})
    if isinstance(result, tuple):
        text_result = result[0]
    elif hasattr(result, "content"):
        text_result = result.content
    else:
        text_result = str(result)

    slides_content = text_result.split("\n## ")
    
    # Generate and store teacher-like explanations for each slide
    extra_explanations = {}
    for i, slide in enumerate(slides_content):
        if slide.strip():
            slide_key = f"Slide {i+1}"
            extra_explanations[slide_key] = generate_extra_explanation(slide, context="topic", lang_choice=lang_choice)
    
    return slides_content, extra_explanations

def generate_slides_from_document(text, instructions=None, num_slides=10, lang_choice='english'):
    """Generate slides from document content with clear, teacher-style explanations, using Markdown formatting."""
    system_prompt = f"""You are a helpful assistant that generates a {num_slides}-slide presentation based on the following document text.
Each slide should be well-structured, concise, and written in a teacher-like explanatory tone.
Ensure a logical progression from introduction to conclusion.
Do not include or mention any images or visual elements, as the output is purely text-based.
{instructions or 'Use standard presentation format with bullet points or short paragraphs for clarity.'}
The output should be in Markdown format, with each slide starting as '## Slide [number]:' followed by its content."""
    # system_prompt is now ready to use in your LLM call

    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("user", "Document Text: {text}"),
    ])
    chain = prompt | llm
    slides_content = chain.invoke({"text": text}).content.split("\n## ")
    
    # Generate and store teacher-like explanations for each slide
    extra_explanations = {}
    for i, slide in enumerate(slides_content):
        if slide.strip():
            slide_key = f"Slide {i+1}"
            extra_explanations[slide_key] = generate_extra_explanation(slide, context="document", lang_choice=lang_choice)
    
    return slides_content, extra_explanations

def create_ppt(slides, filename_prefix="presentation"):
    """Create PowerPoint file with unique filename, each slide on a new page."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{filename_prefix}_{timestamp}.pptx"
    prs = Presentation()
    
    for slide_content in slides:
        if slide_content.strip():
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            lines = slide_content.strip().split("\n")
            title_text = lines[0].replace("## ", "").strip() if lines else "Slide"
            content_text = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
            
            title = new_slide.shapes.title
            title.text = title_text
            content = new_slide.shapes.placeholders[1].text_frame
            content.text = content_text
    
    prs.save(filename)
    print(f"PPT saved as '{filename}'.")
    return filename

def qa_session(topic_slides, extra_explanations, index, chunks, topic, context_type):
    """Question-answer session with teacher-like explanations."""
    print("\n=== Ask Your Questions ===")
    print("Type in English or Roman Urdu.")
    
    while True:
        lang_choice = input("\nAnswer in English or Roman Urdu? (english/roman): ").strip().lower()
        if lang_choice not in ['english', 'roman', 'roman urdu']:
            print("Please choose 'english' or 'roman'.")
            continue
        response_lang = "Roman Urdu" if lang_choice in ['roman', 'roman urdu'] else "English"

        question = input("\nAsk a question (or type 'exit'): ").strip().lower()
        
        if not question:
            continue
        if question in ['exit']:
            break

        system_prompt = f"""You are an experienced teacher. Answer accurately in {response_lang}, matching the question's style. 
        Use the provided context and extra explanations if available, and ensure the response is clear, concise, and engaging with real-life examples."""

        if context_type == 'topic' and topic_slides and extra_explanations:
            relevant_slides = []
            for i, slide in enumerate(topic_slides):
                if question.lower() in slide.lower():
                    relevant_slides.append((slide, extra_explanations.get(f"Slide {i+1}", "")))
            
            if relevant_slides:
                context = "\n".join([f"Slide Content: {slide}\nTeacher Explanation: {exp}" for slide, exp in relevant_slides])
            else:
                context = "\n".join([f"Slide {i+1}: {slide}\nTeacher Explanation: {extra_explanations.get(f'Slide {i+1}', '')}" for i, slide in enumerate(topic_slides)])
            
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"{system_prompt}\nContext: {context}"),
                ("user", "Question: {question}"),
            ])
            chain = prompt | llm
            answer = chain.invoke({"question": question}).content
            print(f"\nAnswer: {answer}")

        elif context_type == 'document' and index and chunks:
            relevant_chunks = retrieve_relevant_chunks(question, index, chunks)
            if relevant_chunks:
                # Generate teacher-like explanation for the relevant document chunks
                doc_explanation = generate_extra_explanation("\n".join(relevant_chunks), context="document", lang_choice=lang_choice)
                prompt = ChatPromptTemplate.from_messages([
                    ("system", f"{system_prompt}\nContext: {relevant_chunks}\nTeacher Explanation: {doc_explanation}"),
                    ("user", "Question: {question}"),
                ])
                chain = prompt | llm
                answer = chain.invoke({"question": question}).content
                print(f"\nAnswer: {answer}")
            else:
                answer = "No relevant info found in the document." if response_lang == "English" else "Document mein is sawaal ka jawab nahi mila."
                print(f"\nAnswer: {answer}")

        else:
            prompt = ChatPromptTemplate.from_messages([
                ("system", system_prompt),
                ("user", "Question: {question}"),
            ])
            chain = prompt | llm
            answer = chain.invoke({"question": question}).content
            print(f"\nAnswer: {answer}")

def main():
    """Main chatbot interface."""
    print("=== Smart Learn AI Chatbot ===")
    print("Use English or Roman Urdu!")
    
    while True:
        print("\n=== Main Menu ===")
        print("1. Enter Topic")
        print("2. Enter Document Path")
        print("3. Exit")
        
        choice = input("\nChoose an option (1-3): ").strip()
        
        if choice == '1':
            topic = input("\nWhat's the presentation topic? ").strip()
            if not topic:
                print("No topic entered, back to menu.")
                continue
            
            instructions = input("Any specific instructions? (e.g., 'Use simple language') or 'skip': ").strip()
            instructions = None if instructions.lower() in ['skip', 'skip karo'] else instructions
            
            num_slides = input("How many slides? (Default 10): ").strip()
            num_slides = int(num_slides) if num_slides.isdigit() else 10
            
            lang_choice = input("Language for explanations (english/roman): ").strip().lower()
            lang_choice = 'english' if lang_choice not in ['roman', 'roman urdu'] else 'roman'
            
            print("\nGenerating topic slides and teacher explanations...")
            topic_slides, extra_explanations = generate_slides_from_topic(topic, instructions, num_slides, lang_choice)
            create_ppt(topic_slides, filename_prefix="topic_presentation")
            
            print("\nGenerating quiz based on topic slides...")
            quiz_content, quiz_filename = generate_quiz(topic_slides, context_type='topic', lang_choice=lang_choice)
            print(f"\n=== Quiz ===\n{quiz_content}")
            
            # Print teacher-like explanations as if presenting
            print("\n=== Teacher's Presentation ===")
            for i, slide in enumerate(topic_slides):
                slide_key = f"Slide {i+1}"
                if slide.strip():
                    print(f"\n{slide}")
                    print(f"Teacher's Explanation:\n{extra_explanations.get(slide_key, 'No explanation generated.')}")
            
            qa_session(topic_slides, extra_explanations, None, None, topic, context_type='topic')
        
        elif choice == '2':
            file_path = input("\nEnter document path (PDF, TXT): ").strip()
            try:
                text = extract_text_from_file(file_path)
                chunks = split_text_into_chunks(text)
                embeddings_array = generate_embeddings(chunks)
                index = create_faiss_index(embeddings_array)
                
                instructions = input("Any specific instructions? (e.g., 'Use simple language') or 'skip': ").strip()
                instructions = None if instructions.lower() in ['skip', 'skip karo'] else instructions
                
                num_slides = input("How many slides? (Default 10): ").strip()
                num_slides = int(num_slides) if num_slides.isdigit() else 10
                
                lang_choice = input("Language for explanations (english/roman): ").strip().lower()
                lang_choice = 'english' if lang_choice not in ['roman', 'roman urdu'] else 'roman'
                
                print("\nGenerating document slides and teacher explanations...")
                doc_slides, extra_explanations = generate_slides_from_document(text, instructions, num_slides, lang_choice)
                create_ppt(doc_slides, filename_prefix="document_presentation")
                
                print("\nGenerating quiz based on document slides...")
                quiz_content, quiz_filename = generate_quiz(doc_slides, context_type='document', lang_choice=lang_choice)
                print(f"\n=== Quiz ===\n{quiz_content}")
                
                # Print teacher-like explanations as if presenting
                print("\n=== Teacher's Presentation ===")
                for i, slide in enumerate(doc_slides):
                    slide_key = f"Slide {i+1}"
                    if slide.strip():
                        print(f"\n{slide}")
                        print(f"Teacher's Explanation:\n{extra_explanations.get(slide_key, 'No explanation generated.')}")
                
                qa_session(None, extra_explanations, index, chunks, None, context_type='document')
                
            except Exception as e:
                print(f"Error with document: {e}")
                continue
        
        elif choice == '3':
            print("Goodbye!")
            break
        
        else:
            print("Invalid choice, pick 1, 2, or 3.")

if __name__ == "__main__":
    main()